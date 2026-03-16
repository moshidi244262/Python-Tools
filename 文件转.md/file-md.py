# 依赖安装: 
# 1. 基础依赖: pip install python-docx PyPDF2 python-pptx pandas tabulate openpyxl tkinterdnd2 pdfplumber
# 2. 支持 .doc 和 .ppt (旧格式) 需要安装: pip install pywin32 (仅限 Windows 系统，且需安装 Microsoft Office/WPS)

import os
import sys
import threading
import platform
import subprocess
from typing import List, Optional

# --- 导入第三方库 ---
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, ttk
from tkinter.constants import END

# 拖拽
try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    DND_SUPPORT = True
except ImportError:
    DND_SUPPORT = False
    class TkinterDnD:
        @staticmethod
        def Tk(): return tk.Tk()

# Word
try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

# PDF
try:
    import pdfplumber
    PDFPLUMBER_SUPPORT = True
except ImportError:
    PDFPLUMBER_SUPPORT = False

try:
    from PyPDF2 import PdfReader
    PYPDF2_SUPPORT = True
except ImportError:
    PYPDF2_SUPPORT = False

# PPT
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

# Excel
try:
    import pandas as pd
    import tabulate
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False

# Windows COM
IS_WINDOWS = platform.system() == 'Windows'
WIN32COM_SUPPORT = False
if IS_WINDOWS:
    try:
        import win32com.client
        import pythoncom  # 多线程必须
        WIN32COM_SUPPORT = True
    except ImportError:
        pass


class FileToMarkdownConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("通用文件转 Markdown 工具 Pro v7.0")
        self.root.geometry("1100x750")
        
        # UI 字体与样式配置
        self.default_font = ("Microsoft YaHei", 10)
        self.root.option_add("*Font", self.default_font)
        self.style = ttk.Style()
        if IS_WINDOWS:
            self.style.theme_use('vista')
            
        # 路径配置
        if getattr(sys, 'frozen', False):
            self.script_dir = os.path.dirname(sys.executable)
        else:
            self.script_dir = os.path.dirname(os.path.abspath(__file__))
        self.output_base_dir = os.path.join(self.script_dir, "GLM-md")
        
        # 支持的格式
        self.supported_ext = {
            '.doc', '.docx', '.pdf', '.ppt', '.pptx', '.rtf',
            '.txt', '.md', '.py', '.js', '.java', '.c', '.cpp', '.h', '.hpp', 
            '.cs', '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.scala', 
            '.html', '.css', '.scss', '.xml', '.json', '.yaml', '.yml', 
            '.ini', '.cfg', '.conf', '.sh', '.bat', '.cmd', '.log', '.sql',
            '.csv', '.xlsx', '.xls'
        }
        
        self.task_list = []
        self.is_converting = False
        self.stop_event = threading.Event()  # 用于中止转换
        
        self.setup_ui()

    def setup_ui(self):
        main_frame = tk.Frame(self.root, bg="#f5f5f7")
        main_frame.pack(fill=tk.BOTH, expand=True, padx=15, pady=15)
        
        # --- 左侧 (操作区) ---
        left_frame = tk.Frame(main_frame, bg="#f5f5f7")
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        # 1. 工具栏
        toolbar = tk.Frame(left_frame, bg="#f5f5f7", pady=5)
        toolbar.pack(fill=tk.X)
        
        self.create_flat_button(toolbar, "添加文件", self.select_files, "#e0e0e0").pack(side=tk.LEFT, padx=3)
        self.create_flat_button(toolbar, "添加文件夹", self.select_folder, "#e0e0e0").pack(side=tk.LEFT, padx=3)
        self.btn_convert = self.create_flat_button(toolbar, "▶ 开始转换", self.start_conversion_thread, "#4CAF50", "white")
        self.btn_convert.pack(side=tk.LEFT, padx=10)
        self.btn_stop = self.create_flat_button(toolbar, "■ 停止", self.stop_conversion, "#F44336", "white")
        self.btn_stop.pack(side=tk.LEFT, padx=3)
        self.btn_stop.config(state=tk.DISABLED)

        # 2. 列表管理区
        btn_frame_2 = tk.Frame(left_frame, bg="#f5f5f7")
        btn_frame_2.pack(fill=tk.X, pady=8)
        
        self.create_flat_button(btn_frame_2, "移除选中项", self.remove_selected_items, "#ffcccc").pack(side=tk.LEFT, padx=3)
        self.create_flat_button(btn_frame_2, "清空列表", self.clear_list, "#ffe0b2").pack(side=tk.LEFT, padx=3)
        self.create_flat_button(btn_frame_2, "清空日志", self.clear_log, "#e0e0e0").pack(side=tk.RIGHT, padx=3)

        # 3. 待转换列表 (支持拖拽)
        list_label_frame = tk.Frame(left_frame, bg="#f5f5f7")
        list_label_frame.pack(fill=tk.BOTH, expand=True)
        
        tk.Label(list_label_frame, text="待转换列表 (支持拖拽文件/文件夹至此处):", bg="#f5f5f7", font=("Microsoft YaHei", 9, "bold")).pack(anchor=tk.W, pady=(0, 5))
        
        self.listbox = tk.Listbox(list_label_frame, selectmode=tk.EXTENDED, bg="#ffffff", relief=tk.FLAT, highlightthickness=1, highlightcolor="#0078D7", font=("Consolas", 10))
        self.listbox.pack(fill=tk.BOTH, expand=True, side=tk.LEFT)
        
        scrollbar = ttk.Scrollbar(list_label_frame, command=self.listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.listbox.config(yscrollcommand=scrollbar.set)
        
        if DND_SUPPORT:
            self.listbox.drop_target_register(DND_FILES)
            self.listbox.dnd_bind('<<Drop>>', self.on_drop)
            self.listbox.dnd_bind('<<DragEnter>>', lambda e: self.listbox.config(bg="#e8f4fa"))
            self.listbox.dnd_bind('<<DragLeave>>', lambda e: self.listbox.config(bg="#ffffff"))
        
        # 4. 进度条与状态
        progress_frame = tk.Frame(left_frame, bg="#f5f5f7", pady=10)
        progress_frame.pack(fill=tk.X)
        
        self.status_label = tk.Label(progress_frame, text="就绪", bg="#f5f5f7", fg="#555555")
        self.status_label.pack(side=tk.LEFT)
        
        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(progress_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.pack(side=tk.RIGHT, fill=tk.X, expand=True, padx=(10, 0))

        # 5. 日志与输出区
        log_frame = tk.Frame(left_frame, bg="#f5f5f7")
        log_frame.pack(fill=tk.BOTH, expand=True)
        
        path_frame = tk.Frame(log_frame, bg="#f5f5f7")
        path_frame.pack(fill=tk.X, pady=(0, 5))
        tk.Label(path_frame, text=f"输出目录: {self.output_base_dir}", fg="#005a9e", bg="#f5f5f7", font=("Microsoft YaHei", 9, "bold")).pack(side=tk.LEFT)
        self.create_flat_button(path_frame, "打开目录", self.open_output_dir, "#b3d4ff").pack(side=tk.RIGHT)
        
        self.log_area = scrolledtext.ScrolledText(log_frame, height=8, state='disabled', bg="#1e1e1e", fg="#cccccc", font=("Consolas", 9), relief=tk.FLAT)
        self.log_area.pack(fill=tk.BOTH, expand=True)
        
        # 配置日志颜色标签
        self.log_area.tag_config("info", foreground="#cccccc")
        self.log_area.tag_config("success", foreground="#4CAF50")
        self.log_area.tag_config("error", foreground="#F44336")
        self.log_area.tag_config("warning", foreground="#FFEB3B")

        # --- 右侧 (信息区) ---
        right_frame = tk.Frame(main_frame, bg="#f5f5f7", width=250)
        right_frame.pack(side=tk.RIGHT, fill=tk.Y, padx=(15, 0))
        
        info_label = tk.LabelFrame(right_frame, text="运行环境与支持格式", bg="#ffffff", font=("Microsoft YaHei", 9, "bold"), padx=10, pady=10, relief=tk.FLAT)
        info_label.pack(fill=tk.BOTH, expand=True)
        
        info_text = tk.Text(info_label, width=28, wrap=tk.WORD, bg="#ffffff", relief=tk.FLAT, font=("Microsoft YaHei", 9))
        info_text.pack(fill=tk.BOTH, expand=True)
        
        # 插入环境信息
        info_text.insert(END, "【引擎状态】\n", "bold")
        info_text.insert(END, f"PDF: {'pdfplumber' if PDFPLUMBER_SUPPORT else ('PyPDF2' if PYPDF2_SUPPORT else '未安装')}\n", "ok" if PDFPLUMBER_SUPPORT else ("warn" if PYPDF2_SUPPORT else "err"))
        info_text.insert(END, f"Office(Win): {'可用' if WIN32COM_SUPPORT else '未安装'}\n", "ok" if WIN32COM_SUPPORT else "err")
        info_text.insert(END, f"Excel: {'可用' if EXCEL_SUPPORT else '未安装'}\n", "ok" if EXCEL_SUPPORT else "err")
        info_text.insert(END, f"拖拽: {'可用' if DND_SUPPORT else '未安装'}\n\n", "ok" if DND_SUPPORT else "err")
        
        info_text.insert(END, "【文档类】\n.doc, .docx, .pdf, .ppt, .pptx, .rtf\n\n", "normal")
        info_text.insert(END, "【数据类】\n.csv, .xlsx, .xls, .json, .xml, .yaml\n\n", "normal")
        info_text.insert(END, "【代码文本类】\n.py, .java, .c, .cpp, .html, .css, .txt, .md ...\n", "normal")
        
        info_text.tag_config("bold", font=("Microsoft YaHei", 9, "bold"))
        info_text.tag_config("ok", foreground="#2e7d32")
        info_text.tag_config("warn", foreground="#f57f17")
        info_text.tag_config("err", foreground="#c62828")
        info_text.tag_config("normal", foreground="#444444")
        info_text.config(state='disabled')

    def create_flat_button(self, parent, text, command, bg, fg="black"):
        """创建扁平化按钮"""
        btn = tk.Button(parent, text=text, command=command, bg=bg, fg=fg, relief=tk.FLAT, bd=0, padx=10, pady=3, cursor="hand2", font=("Microsoft YaHei", 9))
        # 简单的悬浮变色反馈
        def on_enter(e): e.widget['background'] = self.adjust_color(bg, -20)
        def on_leave(e): e.widget['background'] = bg
        btn.bind("<Enter>", on_enter)
        btn.bind("<Leave>", on_leave)
        return btn

    def adjust_color(self, hex_color, factor):
        """简单调整十六进制颜色的亮度用于悬浮效果"""
        hex_color = hex_color.lstrip('#')
        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        new_rgb = tuple(max(0, min(255, c + factor)) for c in rgb)
        return f"#{new_rgb[0]:02x}{new_rgb[1]:02x}{new_rgb[2]:02x}"

    def log(self, message: str, level="info"):
        self.root.after(0, self._log_safe, message, level)

    def _log_safe(self, message: str, level: str):
        self.log_area.config(state='normal')
        self.log_area.insert(tk.END, message + "\n", level)
        self.log_area.see(tk.END)
        self.log_area.config(state='disabled')

    def clear_log(self):
        self.log_area.config(state='normal')
        self.log_area.delete(1.0, END)
        self.log_area.config(state='disabled')

    def clear_list(self):
        self.task_list.clear()
        self.listbox.delete(0, END)
        self.log("列表已清空。")

    def open_output_dir(self):
        if not os.path.exists(self.output_base_dir):
            os.makedirs(self.output_base_dir)
        if IS_WINDOWS:
            os.startfile(self.output_base_dir)
        elif platform.system() == "Darwin":  # macOS
            subprocess.Popen(["open", self.output_base_dir])
        else: # Linux
            subprocess.Popen(["xdg-open", self.output_base_dir])

    def update_status(self, text: str, progress: float = None):
        def _update():
            self.status_label.config(text=text)
            if progress is not None:
                self.progress_var.set(progress)
        self.root.after(0, _update)

    def _add_paths_to_queue(self, paths: List[str], base_source_dir: str = None):
        added_count = 0
        for path in paths:
            path = path.strip('{}')
            if not os.path.exists(path): continue

            if os.path.isfile(path):
                _, ext = os.path.splitext(path)
                if ext.lower() in self.supported_ext:
                    rel_path = os.path.relpath(path, base_source_dir) if base_source_dir else os.path.basename(path)
                    if path not in [t['src'] for t in self.task_list]:
                        self.task_list.append({'src': path, 'dst_rel': rel_path})
                        added_count += 1
                else:
                    self.log(f"跳过不支持的格式: {os.path.basename(path)}", "warning")

            elif os.path.isdir(path):
                current_base = path
                for root, dirs, files in os.walk(path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        _, ext = os.path.splitext(full_path)
                        if ext.lower() in self.supported_ext:
                            rel_path = os.path.relpath(full_path, current_base)
                            if full_path not in [t['src'] for t in self.task_list]:
                                self.task_list.append({'src': full_path, 'dst_rel': rel_path, 'base': current_base})
                                added_count += 1

        self.refresh_listbox()
        return added_count

    def refresh_listbox(self):
        self.listbox.delete(0, END)
        for task in self.task_list:
            display_text = f" {os.path.basename(task['src'])}  |  {task['dst_rel']}"
            self.listbox.insert(END, display_text)

    def select_files(self):
        files = filedialog.askopenfilenames(title="选择要转换的文件")
        if files:
            count = self._add_paths_to_queue(files)
            self.log(f"添加了 {count} 个文件。")

    def select_folder(self):
        folder = filedialog.askdirectory(title="选择要转换的文件夹")
        if folder:
            count = self._add_paths_to_queue([folder], base_source_dir=folder)
            self.log(f"从文件夹添加了 {count} 个文件。")

    def on_drop(self, event):
        raw_data = event.data
        paths = []
        if '{' in raw_data:
            temp = ""
            in_brace = False
            for char in raw_data:
                if char == '{': in_brace = True
                elif char == '}': 
                    in_brace = False
                    paths.append(temp)
                    temp = ""
                elif in_brace: temp += char
                elif char == ' ' and not in_brace:
                    if temp: paths.append(temp); temp = ""
                else: temp += char
            if temp: paths.append(temp)
        else:
            paths = raw_data.split()

        if paths:
            count = 0
            for p in paths:
                count += self._add_paths_to_queue([p], base_source_dir=p if os.path.isdir(p) else None)
            self.log(f"拖拽添加了 {count} 个项目。")
            self.listbox.config(bg="#ffffff")

    def remove_selected_items(self):
        selection = self.listbox.curselection()
        if not selection: return
        for i in sorted(selection, reverse=True):
            del self.task_list[i]
            self.listbox.delete(i)

    def stop_conversion(self):
        if self.is_converting:
            self.stop_event.set()
            self.update_status("正在中止...", self.progress_var.get())
            self.btn_stop.config(state=tk.DISABLED)

    def start_conversion_thread(self):
        if self.is_converting: return
        if not self.task_list:
            messagebox.showwarning("提示", "列表为空，请先添加文件！")
            return
        
        self.is_converting = True
        self.stop_event.clear()
        self.btn_convert.config(state=tk.DISABLED)
        self.btn_stop.config(state=tk.NORMAL)
        self.progress_var.set(0)
        
        thread = threading.Thread(target=self.convert_files_thread, daemon=True)
        thread.start()

    def convert_files_thread(self):
        # 多线程环境下必须初始化 COM
        if WIN32COM_SUPPORT:
            pythoncom.CoInitialize()
            
        word_app, ppt_app = None, None
        if WIN32COM_SUPPORT:
            try:
                # 预启动进程以复用，提升速度
                word_app = win32com.client.DispatchEx("Word.Application")
                word_app.Visible = False
                ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
            except Exception as e:
                self.log(f"COM 初始化警告 (部分老格式可能失败): {e}", "warning")

        try:
            if not os.path.exists(self.output_base_dir):
                os.makedirs(self.output_base_dir)
        except Exception as e:
            self.log(f"错误: 无法创建输出文件夹 - {e}", "error")
            self.on_conversion_finished()
            return

        success_count = 0
        fail_count = 0
        total = len(self.task_list)
        
        for idx, task in enumerate(self.task_list):
            if self.stop_event.is_set():
                self.log("用户已强行中止转换任务。", "warning")
                break
                
            src_path = task['src']
            dst_rel = task['dst_rel']
            file_name = os.path.basename(src_path)
            
            progress_pct = (idx / total) * 100
            self.update_status(f"处理中 ({idx+1}/{total}): {file_name}", progress_pct)
            self.log(f"正在转换: {file_name} ...")

            try:
                md_content = self.process_single_file(src_path, word_app, ppt_app)
                
                if md_content is not None:
                    dst_rel_md = os.path.splitext(dst_rel)[0] + ".md"
                    final_save_path = os.path.join(self.output_base_dir, dst_rel_md)
                    
                    os.makedirs(os.path.dirname(final_save_path), exist_ok=True)
                        
                    with open(final_save_path, 'w', encoding='utf-8') as f:
                        f.write(md_content)
                    
                    self.log(f"  成功 -> {dst_rel_md}", "success")
                    success_count += 1
                else:
                    self.log(f"  跳过: 不支持或为空文件", "warning")
                    fail_count += 1
                        
            except Exception as e:
                self.log(f"  失败: {file_name} - {str(e)}", "error")
                fail_count += 1

        # 清理 COM 进程释放内存
        if word_app: 
            try: word_app.Quit()
            except: pass
        if ppt_app: 
            try: ppt_app.Quit()
            except: pass
            
        if WIN32COM_SUPPORT:
            pythoncom.CoUninitialize()

        self.update_status("完成" if not self.stop_event.is_set() else "已中止", 100)
        self.log(f"\n--- 任务结束 | 成功: {success_count} | 失败/跳过: {fail_count} ---", "info")
        
        if not self.stop_event.is_set():
            self.root.after(0, lambda: messagebox.showinfo("完成", f"转换完成！\n成功: {success_count}\n失败/跳过: {fail_count}"))
        
        self.on_conversion_finished()

    def on_conversion_finished(self):
        self.is_converting = False
        self.stop_event.clear()
        self.root.after(0, lambda: self.btn_convert.config(state=tk.NORMAL))
        self.root.after(0, lambda: self.btn_stop.config(state=tk.DISABLED))


    # ============== 文件解析核心逻辑 ==============

    def process_single_file(self, file_path: str, word_app=None, ppt_app=None) -> Optional[str]:
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext not in self.supported_ext: return None

        if ext == '.doc': return self.convert_doc_win32(file_path, word_app)
        elif ext == '.ppt': return self.convert_ppt_win32(file_path, ppt_app)
        elif ext == '.rtf': return self.convert_rtf_win32(file_path, word_app)

        if ext == '.docx': return self.convert_docx_optimized(file_path)
        elif ext == '.pdf': return self.convert_pdf_optimized(file_path)
        elif ext == '.pptx': return self.convert_pptx(file_path)
        elif ext in ('.xlsx', '.xls'): return self.convert_excel(file_path)
        elif ext == '.csv': return self.convert_csv(file_path)
        else: return self.convert_text_or_code(file_path, ext)

    def convert_docx_optimized(self, file_path: str) -> str:
        if not DOCX_SUPPORT: raise Exception("未安装 python-docx")
        doc = Document(file_path)
        md_content = f"# {os.path.basename(file_path)}\n\n"
        
        from docx.oxml.ns import qn
        
        def is_inside_table(paragraph):
            return paragraph._element.getparent().tag == qn('w:tc') 

        para_map = {p._element: p for p in doc.paragraphs}
        table_map = {t._tbl: t for t in doc.tables}
        
        for element in doc.element.body.iterchildren():
            if element in para_map:
                p = para_map[element]
                if is_inside_table(p): continue
                text = p.text.strip()
                if not text: continue
                style_name = p.style.name.lower() if p.style else ""
                
                if 'heading 1' in style_name or 'heading1' in style_name: md_content += f"# {text}\n\n"
                elif 'heading 2' in style_name or 'heading2' in style_name: md_content += f"## {text}\n\n"
                elif 'heading 3' in style_name or 'heading3' in style_name: md_content += f"### {text}\n\n"
                else: md_content += f"{text}\n\n"
                    
            elif element in table_map:
                table = table_map[element]
                md_content += self._convert_table_to_md(table) + "\n"

        return md_content

    def _convert_table_to_md(self, table) -> str:
        md_content = ""
        try:
            col_count = len(table.columns) if len(table.columns) > 0 else (len(table.rows[0].cells) if len(table.rows) > 0 else 0)
            if col_count == 0: return ""
            
            temp_content = []
            for r_idx in range(len(table.rows)):
                cols_in_row = []
                for c_idx in range(col_count):
                    try: cell = table.cell(r_idx, c_idx)
                    except Exception: cols_in_row.append(" "); continue
                    
                    is_v_merge, is_h_merge = False, False
                    if r_idx > 0:
                        try:
                            if cell is table.cell(r_idx - 1, c_idx): is_v_merge = True
                        except: pass
                    if c_idx > 0:
                        try:
                            if cell is table.cell(r_idx, c_idx - 1): is_h_merge = True
                        except: pass
                    
                    if is_v_merge or is_h_merge:
                        cols_in_row.append(" ")
                    else:
                        text = cell.text.strip().replace('\n', '<br>').replace('|', '&#124;')
                        cols_in_row.append(text if text else " ")
                        
                temp_content.append(cols_in_row)
                
            for r_idx, row_vals in enumerate(temp_content):
                md_content += "| " + " | ".join(row_vals) + " |\n"
                if r_idx == 0:
                    md_content += "| " + " | ".join(["---"] * len(row_vals)) + " |\n"
        except Exception:
            md_content += "\n[表格解析遇到复杂结构，已降级处理]\n"
            for row in table.rows:
                cells_text = [cell.text.strip().replace('\n', '<br>') for cell in row.cells]
                md_content += "| " + " | ".join(cells_text) + " |\n"
        return md_content

    def convert_pdf_optimized(self, file_path: str) -> str:
        md_content = f"# {os.path.basename(file_path)}\n\n"
        if PDFPLUMBER_SUPPORT:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text: md_content += f"## Page {i+1}\n\n{text}\n\n"
                        for table in page.extract_tables():
                            for row in table:
                                clean_row = [str(cell).strip().replace('\n', ' ') if cell else " " for cell in row]
                                md_content += "| " + " | ".join(clean_row) + " |\n"
                            md_content += "\n"
                return md_content
            except Exception as e:
                self.log(f"pdfplumber 解析失败: {e}，尝试 PyPDF2...", "warning")

        if PYPDF2_SUPPORT:
            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text: md_content += f"## Page {i+1}\n\n{text}\n\n"
            return md_content
        raise Exception("未安装 PDF 解析库 (请安装 pdfplumber)")

    def convert_pptx(self, file_path: str) -> str:
        if not PPTX_SUPPORT: raise Exception("未安装 python-pptx")
        prs = Presentation(file_path)
        md_content = f"# {os.path.basename(file_path)}\n\n"
        for i, slide in enumerate(prs.slides):
            md_content += f"---\n## 幻灯片 {i+1}\n\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    md_content += f"{shape.text.strip()}\n\n"
        return md_content

    def convert_excel(self, file_path: str) -> str:
        if not EXCEL_SUPPORT: raise Exception("未安装 pandas / tabulate")
        xlsx = pd.ExcelFile(file_path)
        md_content = f"# {os.path.basename(file_path)}\n\n"
        for sheet_name in xlsx.sheet_names:
            df = pd.read_excel(xlsx, sheet_name=sheet_name)
            md_table = tabulate.tabulate(df, headers='keys', tablefmt='pipe', showindex=False)
            md_content += f"## Sheet: {sheet_name}\n\n{md_table}\n\n"
        return md_content

    def convert_csv(self, file_path: str) -> str:
        if not EXCEL_SUPPORT:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f"# {os.path.basename(file_path)}\n\n```csv\n{f.read()}\n```"
        df = pd.read_csv(file_path)
        md_table = tabulate.tabulate(df, headers='keys', tablefmt='pipe', showindex=False)
        return f"# {os.path.basename(file_path)}\n\n{md_table}"

    def convert_text_or_code(self, file_path: str, ext: str) -> str:
        content = ""
        for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
            try:
                with open(file_path, 'r', encoding=enc) as f: 
                    content = f.read(); break
            except: continue
        
        if not content: return f"# {os.path.basename(file_path)}\n\n*文件内容为空或无法解码*"
        
        header = f"# {os.path.basename(file_path)}\n\n"
        code_exts = {'.py', '.js', '.java', '.c', '.cpp', '.h', '.cs', '.go', '.rs', '.rb', 
                     '.php', '.swift', '.kt', '.scala', '.html', '.css', '.xml', '.json', 
                     '.yaml', '.yml', '.sh', '.sql'}
        if ext in code_exts:
            lang = ext[1:] if ext != '.json' else 'json'
            return header + f"```{lang}\n{content}\n```"
        return header + content

    def convert_doc_win32(self, file_path: str, app) -> str:
        if not app: raise Exception("Word COM 初始化失败，无法处理旧版 .doc")
        doc = None
        try:
            doc = app.Documents.Open(os.path.abspath(file_path), ReadOnly=True)
            return f"# {os.path.basename(file_path)}\n\n{doc.Content.Text}\n"
        finally:
            if doc: doc.Close(False)

    def convert_ppt_win32(self, file_path: str, app) -> str:
        if not app: raise Exception("PPT COM 初始化失败，无法处理旧版 .ppt")
        pres = None
        try:
            pres = app.Presentations.Open(os.path.abspath(file_path), ReadOnly=True, Untitled=False, WithWindow=False)
            md_content = f"# {os.path.basename(file_path)}\n\n"
            for i, slide in enumerate(pres.Slides):
                md_content += f"## Slide {i+1}\n\n"
                for shape in slide.Shapes:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        md_content += shape.TextFrame.TextRange.Text + "\n\n"
            return md_content
        finally:
            if pres: pres.Close()

    def convert_rtf_win32(self, file_path: str, app) -> str:
        if not app:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f: return f.read()
        doc = None
        try:
            doc = app.Documents.Open(os.path.abspath(file_path), ReadOnly=True)
            return f"# {os.path.basename(file_path)}\n\n{doc.Content.Text}\n"
        finally:
            if doc: doc.Close(False)


if __name__ == "__main__":
    if DND_SUPPORT:
        root = TkinterDnD.Tk()
    else:
        root = tk.Tk()
    app = FileToMarkdownConverter(root)
    root.mainloop()
